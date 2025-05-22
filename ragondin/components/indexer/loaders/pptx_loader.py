import html
import re
from io import BytesIO
import pptx
from langchain_core.documents.base import Document
from PIL import Image
from tqdm.asyncio import tqdm

from .base import BaseLoader


class PPTXConverter:
    def __init__(
        self, image_placeholder=r"<image>", page_separator: str = "[PAGE_SEP]"
    ):
        self.image_placeholder = image_placeholder
        self.page_separator = page_separator

    def convert(self, local_path):
        md_content = ""
        presentation = pptx.Presentation(local_path)
        slide_num = 0
        images_list = []
        for slide in presentation.slides:
            slide_num += 1

            md_content += self.page_separator

            title = slide.shapes.title
            for shape in slide.shapes:
                if self._is_picture(shape):
                    images_list.append(Image.open(BytesIO(shape.image.blob)))
                    md_content += self.image_placeholder

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return md_content, images_list

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_chart_to_markdown(self, chart):
        md = "\n\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"
        data = []
        category_names = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        data.append(["Category"] + series_names)

        for idx, category in enumerate(category_names):
            row = [category]
            for series in chart.series:
                row.append(series.values[idx])
            data.append(row)

        markdown_table = []
        for row in data:
            markdown_table.append("| " + " | ".join(map(str, row)) + " |")
        header = markdown_table[0]
        separator = "|" + "|".join(["---"] * len(data[0])) + "|"
        return md + "\n".join([header, separator] + markdown_table[1:])


class PPTXLoader(BaseLoader):
    def __init__(self, page_sep: str = "[PAGE_SEP]", **kwargs) -> None:
        super().__init__(page_sep, **kwargs)
        self.image_placeholder = r"<image>"
        self.converter = PPTXConverter(
            image_placeholder=self.image_placeholder, page_separator=page_sep
        )

    async def get_captions(self, images):
        tasks = [self.get_image_description(image=img) for img in images]
        return await tqdm.gather(*tasks, desc="Generating captions")

    async def aload_document(self, file_path, metadata=None, save_markdown=False):
        md_content, imgs = self.converter.convert(local_path=file_path)
        images_captions = await self.get_captions(imgs)

        for caption in images_captions:
            md_content = re.sub(
                self.image_placeholder,
                caption.replace("\\", "/"),
                md_content,
                count=1,
            )

        doc = Document(page_content=md_content, metadata=metadata)
        if save_markdown:
            self.save_document(Document(page_content=md_content), str(file_path))
        return doc
